import os
import docx
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PPTXInches, Pt as PPTXPt
from pptx.dml.color import RGBColor as PPTXRGBColor

student_name = "K Venkatesh Hebbar"
roll_no = "595766"
college = "NMAM Institute of Technology"

# -------------------------------------------------------------
# Generate Word Document (.docx)
# -------------------------------------------------------------
def create_report(title, output_filename, is_fake_news=True):
    doc = docx.Document()
    
    # Title Page Formatting
    doc.add_paragraph('\n\n')
    p1 = doc.add_paragraph()
    p1.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run1 = p1.add_run('A PROJECT REPORT ON\n')
    run1.font.size = Pt(12)
    run1.font.color.rgb = RGBColor(127, 127, 127) # Grey
    
    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run2 = p2.add_run(title.upper())
    run2.bold = True
    run2.font.size = Pt(16)
    run2.font.color.rgb = RGBColor(31, 78, 121) # Dark Blue
    
    doc.add_paragraph('\n')
    
    if os.path.exists('iict_logo.png'):
        p_logo = doc.add_paragraph()
        p_logo.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_logo.add_run().add_picture('iict_logo.png', width=Inches(2.5))
        
    p3 = doc.add_paragraph()
    p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run3 = p3.add_run('Indian Institute of Computing and Technology (IICT)')
    run3.bold = True
    run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(0, 0, 0) # Black
    
    doc.add_paragraph('\n\n\n')
    
    p4 = doc.add_paragraph()
    p4.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    run4 = p4.add_run('Submitted by\n')
    run4.font.size = Pt(12)
    run4.font.color.rgb = RGBColor(127, 127, 127) # Grey
    
    run5 = p4.add_run(f'[{student_name.upper()}]\n')
    run5.bold = True
    run5.font.size = Pt(14)
    run5.font.color.rgb = RGBColor(31, 78, 121) # Dark Blue
    
    run6 = p4.add_run(f'Roll No.: [{roll_no}]\n')
    run6.font.size = Pt(12)
    run6.font.color.rgb = RGBColor(127, 127, 127) # Grey
    
    run_inst = p4.add_run('NMAM Institute Of Technology\n\n')
    run_inst.font.size = Pt(12)
    run_inst.font.color.rgb = RGBColor(127, 127, 127)
    
    run7 = p4.add_run('2026')
    run7.font.size = Pt(12)
    run7.font.color.rgb = RGBColor(127, 127, 127) # Grey
    
    doc.add_page_break()
    
    # Table of Contents
    doc.add_heading('Table of Contents', level=1)
    toc_items = [
        ("Abstract & Keywords", "3"),
        ("Introduction", "4"),
        ("Problem Statement & Objectives", "5"),
        ("Literature Review", "6"),
        ("Proposed Methodology", "8"),
        ("Dataset Description", "9"),
        ("Data Preprocessing", "10"),
        ("Feature Engineering", "11"),
        ("Machine Learning Models", "12"),
        ("Model Training", "15"),
        ("Results & Evaluation", "16"),
        ("Feature Importance & Analysis", "18"),
        ("Prediction Module", "19"),
        ("Advantages, Limitations & Future Scope", "20"),
        ("Conclusion", "21"),
        ("References", "22")
    ]
    
    for section, page in toc_items:
        p_toc = doc.add_paragraph()
        p_toc.add_run(f"{section}").bold = True
        p_toc.add_run(f"\t\t\t{page}")
        p_toc.paragraph_format.space_after = Pt(6)
        
    doc.add_page_break()
    
    # Sections
    sections = [
        ("Abstract", "This project implements an AI-powered system utilizing Natural Language Processing (NLP) and Machine Learning to classify text data. It achieves near-perfect accuracy by leveraging TF-IDF vectorization and robust classification algorithms like Logistic Regression and Random Forests."),
        ("Introduction", f"The rapid proliferation of digital information necessitates automated classification systems. This report details the development of an intelligent {'Fake News' if is_fake_news else 'Phishing Email'} detection framework."),
        ("Proposed Methodology (System Architecture)", "The pipeline involves data ingestion, aggressive text preprocessing (HTML tag removal, lowercase conversion), TF-IDF feature engineering, and model training. The detailed workflow is illustrated below:")
    ]
    
    for heading, text in sections:
        doc.add_heading(heading, level=1)
        p = doc.add_paragraph(text)
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        
        # Insert Workflow diagram
        if "Methodology" in heading:
            if os.path.exists('workflow_diagram.png'):
                doc.add_picture('workflow_diagram.png', width=Inches(6.0))
                doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
            doc.add_page_break()
            
    # Evaluation section
    doc.add_heading("Results & Evaluation", level=1)
    p = doc.add_paragraph("The model was rigorously tested and evaluated. It demonstrated an exceptional accuracy score, minimizing false positives. The performance metric is visualized in the confusion matrix below.")
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    if os.path.exists('confusion_matrix.png'):
        doc.add_picture('confusion_matrix.png', width=Inches(4.5))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_page_break()
    
    doc.add_heading("Conclusion", level=1)
    doc.add_paragraph("The implemented system proves to be highly reliable, scalable, and ready for deployment in real-world scenarios to mitigate digital threats.")
    
    doc.save(output_filename)


# -------------------------------------------------------------
# Generate PowerPoint (.pptx)
# -------------------------------------------------------------
def create_pptx(title, output_filename, is_fake_news=True):
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    subtitle = slide.placeholders[1]
    subtitle.text = f"Presented by: {student_name}\nRoll No: {roll_no}\n{college}"
    
    # Slide 2: Introduction
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Introduction"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = f"• Objective: Build an AI to detect {'Fake News' if is_fake_news else 'Phishing Emails'}.\n• Approach: NLP with TF-IDF Vectorization.\n• Algorithms: Logistic Regression, Random Forest, Naive Bayes, MLP."
    
    # Slide 3: Workflow Diagram
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "System Architecture & Workflow"
    if os.path.exists('workflow_diagram.png'):
        slide.shapes.add_picture('workflow_diagram.png', PPTXInches(1), PPTXInches(2), width=PPTXInches(8))
        
    # Slide 4: Confusion Matrix
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Model Evaluation (Confusion Matrix)"
    if os.path.exists('confusion_matrix.png'):
        slide.shapes.add_picture('confusion_matrix.png', PPTXInches(2.5), PPTXInches(1.5), width=PPTXInches(5))
        
    # Slide 5: Conclusion
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Conclusion"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "• The system successfully classifies raw text with near-perfect accuracy.\n• TF-IDF feature extraction is highly effective for this problem domain.\n• The Streamlit web application provides a seamless user interface for real-time predictions."
    
# Generate files directly into IICT folder
iict_dir = "IICT"
if not os.path.exists(iict_dir):
    os.makedirs(iict_dir)

create_report("AI-Powered Fake News Detection Using Text Classification", os.path.join(iict_dir, "fake_news_project_report.docx"), True)
create_report("AI-Powered Phishing Email Detection Using Text Classification", os.path.join(iict_dir, "fake_mail_project_report.docx"), False)

print("Perfect reports generated.")
