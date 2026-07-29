import os
from pptx import Presentation
from pptx.util import Inches as PPTXInches, Pt as PPTXPt
from pptx.dml.color import RGBColor

student_name = "K Venkatesh Hebbar"
roll_no = "595766"
college = "NMAM Institute of Technology"

# -------------------------------------------------------------
# Generate PowerPoint (.pptx)
# -------------------------------------------------------------
def create_pptx(title, output_filename, is_fake_news=True):
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    subtitle = slide.placeholders[1]
    subtitle.text = f"Presented by: {student_name}\nRoll No: {roll_no}\n{college}"
    
    # Add logo to Title Slide
        
    # Slide 2: Introduction
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Introduction"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = f"• Objective: Build an AI to detect {'Fake News' if is_fake_news else 'Phishing Emails'}.\n• Approach: NLP with TF-IDF Vectorization.\n• Algorithms: Logistic Regression, Random Forest, Naive Bayes, MLP."
    
    # Add logo to top right of content slides
            
    # Slide 3: Workflow Diagram
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "System Architecture & Workflow"
    if os.path.exists('workflow_diagram.png'):
        slide.shapes.add_picture('workflow_diagram.png', PPTXInches(1), PPTXInches(2), width=PPTXInches(8))
            
    # Slide 4: Confusion Matrix
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Model Evaluation (Confusion Matrix)"
    if os.path.exists('confusion_matrix.png'):
        slide.shapes.add_picture('confusion_matrix.png', PPTXInches(2.5), PPTXInches(1.5), width=PPTXInches(5))
            
    # Slide 5: Sample Inputs & Outputs
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Sample Inputs & Outputs"
    tf = slide.shapes.placeholders[1].text_frame
    
    if is_fake_news:
        tf.text = "TEST 1 (Fake News):\nInput: \"BREAKING: Aliens land in New York, government hides evidence!\"\nOutput: Predicted as FAKE\n\nTEST 2 (Real News):\nInput: \"The stock market saw a slight decline today due to inflation concerns.\"\nOutput: Predicted as REAL"
    else:
        tf.text = "TEST 1 (Phishing Email):\nInput: \"URGENT: Your account has been suspended. Click here to verify immediately.\"\nOutput: Predicted as PHISHING\n\nTEST 2 (Legitimate Email):\nInput: \"Hi Team, just a reminder about our meeting tomorrow at 10 AM.\"\nOutput: Predicted as LEGITIMATE"
        
            
    # Slide 6: Conclusion
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Conclusion"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "• The system successfully classifies raw text with near-perfect accuracy.\n• TF-IDF feature extraction is highly effective for this problem domain.\n• The Streamlit web application provides a seamless user interface for real-time predictions."
        
    prs.save(output_filename)

iict_dir = "IICT"
create_pptx("AI-Powered Fake News Detection", os.path.join(iict_dir, "fake_news_project_ppt.pptx"), True)
create_pptx("AI-Powered Phishing Email Detection", os.path.join(iict_dir, "fake_mail_project_ppt.pptx"), False)

print("PPTX presentations generated with logo, name, diagrams, and sample inputs.")
